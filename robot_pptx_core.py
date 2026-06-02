# -*- coding: utf-8 -*-
"""
robot_pptx_core.py
ロボット甲子園 アイデアシート → プレゼンPowerPoint 生成コア

SalesNexus AI（旧 名刺管理）の research_core.py に相当する共有ロジック。
Streamlit 非依存。app から import して使う。

ワークフロー:
  1) ocr_idea_sheet()        手書きアイデアシートを Gemini 2.5 Flash で構造化
  2) generate_presentation() OCR結果を Ollama gpt-oss:120b-cloud でスライド原稿化
  3) build_pptx()            python-pptx で白背景・非ビジネス調のスライド生成

依存: google-generativeai, ollama, python-pptx, pymupdf(任意), Pillow
"""

import io
import os
import json
import re
import hashlib
from datetime import datetime

# ----------------------------------------------------------------------------
# 定数
# ----------------------------------------------------------------------------
# コスト/濫用対策のサーバー側上限（クライアントからは変更不可）
MAX_INPUT_CHARS = 8000        # 生成に渡す入力テキストの最大文字数
MAX_OUTPUT_TOKENS = 2048      # 生成の最大出力トークン

OCR_MODEL = "gemini-2.5-flash"            # 手書きOCR（クラウド・高精度・安価）
OCR_VISION_MODEL = "qwen3-vl:235b-cloud"  # 手書きOCR代替（Ollamaビジョン。ローカルは qwen3-vl:8b）
GEN_MODEL_DEFAULT = "gpt-oss:20b-cloud"    # 原稿生成（軽量級・負荷レベル1。重いなら120bに変更可）

# 配色（学生向け・元気だが白背景を生かしたクリーンな配色）
COL_INK = "1F2937"      # 本文（濃いグレー）
COL_MUTE = "6B7280"     # 補足
COL_BG_CARD = "F3F4F6"  # 中に置く薄いカード（※フチには触れさせない）
COL_A1 = "EA5A2D"       # アクセント1（オレンジ）
COL_A2 = "2563EB"       # アクセント2（ブルー）
COL_A3 = "16A34A"       # アクセント3（グリーン）
ACCENTS = [COL_A1, COL_A2, COL_A3]

# Windows標準の親しみやすい日本語フォント
FONT_HEAD = "Meiryo"
FONT_BODY = "Meiryo"


# ============================================================================
# 1. OCR（Gemini 2.5 Flash）
# ============================================================================
def _img_bytes_to_part(image_bytes, mime="image/png"):
    return {"mime_type": mime, "data": image_bytes}


OCR_PROMPT = """あなたは手書きアンケート/アイデアシートを読み取る専門OCRです。
画像は、ロボット系イベントで高校生・高専生が手書きしたロボットのアイデアシートです。
すべての手書き文字・図のキャプションを読み取り、次のJSONだけを出力してください（前置き・説明・```不要）。
読み取れない項目は空文字 "" にしてください。推測で創作しないこと。

{
  "robot_name": "ロボットの名称",
  "school": "学校名",
  "grade": "学年・組",
  "author": "氏名",
  "impression": "産業用ロボットを見た感想（あれば）",
  "usage_steps": ["使い方の箇条書きを1項目ずつ"],
  "shape_note": "形・大きさの説明（図のキャプション含む）",
  "motivation": "なぜこの作業をさせようと思ったか",
  "target_demand": "誰が使う・需要に関する記述",
  "price": "想定価格"
}
"""


def ocr_idea_sheet(image_bytes, gemini_api_key, mime="image/png"):
    """手書きアイデアシート画像を Gemini 2.5 Flash で構造化JSONに変換して dict を返す。"""
    import google.generativeai as genai
    genai.configure(api_key=gemini_api_key)
    model = genai.GenerativeModel(OCR_MODEL)
    resp = model.generate_content(
        [OCR_PROMPT, _img_bytes_to_part(image_bytes, mime)],
        generation_config={"temperature": 0.1},
    )
    return _safe_json(resp.text, default=_empty_idea())


def ocr_idea_sheet_ollama(image_bytes, model=OCR_VISION_MODEL, host=None):
    """手書きアイデアシート画像を Ollama のビジョンモデル(qwen3-vl)で構造化JSONに変換。
    Geminiキーが使えないときの代替経路。images にバイト列を渡す。"""
    client = _ollama_client(host)
    resp = client.chat(
        model=model,
        messages=[{"role": "user", "content": OCR_PROMPT, "images": [image_bytes]}],
        options={"temperature": 0.1},
    )
    return _safe_json(resp["message"]["content"], default=_empty_idea())


def run_ocr(image_bytes, *, engine="gemini", mime="image/png",
            gemini_api_key=None, vision_model=OCR_VISION_MODEL, host=None):
    """OCRエンジンのディスパッチャ。engine='gemini' または 'ollama'。"""
    if engine == "ollama":
        return ocr_idea_sheet_ollama(image_bytes, model=vision_model, host=host)
    return ocr_idea_sheet(image_bytes, gemini_api_key, mime=mime)


def _empty_idea():
    return {
        "robot_name": "", "school": "", "grade": "", "author": "",
        "impression": "", "usage_steps": [], "shape_note": "",
        "motivation": "", "target_demand": "", "price": "",
    }


# ============================================================================
# 2. スライド原稿生成（Ollama gpt-oss:120b-cloud）
# ============================================================================
def _ollama_client(host=None):
    """Ollamaクライアントを返す。
    OLLAMA_API_KEY があれば ollama.com に直接接続（ローカルOllama不要・クラウド/サーバー向け）。
    それ以外は OLLAMA_HOST かローカル(localhost:11434)を使用。"""
    import ollama
    api_key = os.environ.get("OLLAMA_API_KEY")
    host = host or os.environ.get("OLLAMA_HOST")
    if api_key:
        return ollama.Client(host=host or "https://ollama.com",
                             headers={"Authorization": "Bearer " + api_key})
    if host:
        return ollama.Client(host=host)
    return ollama


# ============================================================================
# 1.5 Web検索リサーチ（Gemini 2.5 Flash + Google Search グラウンディング）
# ============================================================================
RESEARCH_MODEL = "gemini-2.5-flash"
MAX_RESEARCH_OUTPUT_TOKENS = 3072

RESEARCH_PROMPT = """あなたは中学生・高校生のロボット発表を支援するリサーチャーです。
次のロボットのアイデアについて Google検索を使い、発表で説得力を生む
「事実・統計・背景・関連事例」を箇条書きで集めてください。

# ロボットのアイデア
{idea_text}

# 集める観点
1. 解決しようとしている課題に関する事実・統計（日本国内のデータがあれば優先）
2. 既存の類似ロボット・製品・サービス（中学生にも分かるレベルで2〜4例）
3. このロボットに関連する技術（音声認識、AI、センサーなど）の最新動向
4. ターゲット利用者層の規模・特徴・困りごと
5. 発表で使える「驚きの数字」や「意外な事実」

# 出力フォーマット（厳守）
- 観点ごとに見出しを付け、その下に箇条書き3〜5項目
- 各項目は1〜2文（80〜140字）で具体的に書く。数字は必ず単位とともに
- 各項目末尾に「（出典：URL）」を付ける（検索で得られた場合）
- 見つからない観点は省略してよい
- 事実をねつ造しないこと。推測なら「〜と考えられる」と明記
"""


def gemini_research_topic(idea_data, gemini_api_key):
    """Gemini 2.5 Flash + Google検索でプレゼン用の背景情報を調査する。

    Returns: str（調査結果のテキスト。失敗時は空文字）
    """
    if not gemini_api_key:
        return ""
    import json as _json
    import urllib.request
    import urllib.error

    idea_text = _json.dumps(idea_data, ensure_ascii=False, indent=2)[:MAX_INPUT_CHARS]
    prompt = RESEARCH_PROMPT.format(idea_text=idea_text)
    url = ("https://generativelanguage.googleapis.com/v1beta/models/"
           f"{RESEARCH_MODEL}:generateContent?key={gemini_api_key}")
    body = {
        "contents": [{"parts": [{"text": prompt}]}],
        "tools": [{"google_search": {}}],
        "generationConfig": {
            "temperature": 0.3,
            "maxOutputTokens": MAX_RESEARCH_OUTPUT_TOKENS,
        },
    }
    req = urllib.request.Request(
        url, data=_json.dumps(body).encode("utf-8"),
        headers={"Content-Type": "application/json"}, method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=90) as resp:
            data = _json.loads(resp.read().decode("utf-8"))
        candidates = data.get("candidates", [])
        if not candidates:
            return ""
        parts = candidates[0].get("content", {}).get("parts", [])
        text = "\n".join(p.get("text", "") for p in parts if isinstance(p, dict) and "text" in p)
        return text.strip()
    except Exception:
        # 検索失敗は致命的ではない（リサーチ無しで生成にフォールバック）
        return ""


GEN_SYSTEM = """あなたは、ロボット甲子園に出場する高校生・高専生のアイデアを、
聴衆（同世代の学生・先生・審査員）にワクワクして伝えるプレゼン原稿に仕上げる編集者です。
かたいビジネス文書ではなく、前向きで分かりやすい、はずむような言葉を使ってください。
専門用語は噛み砕き、1行は短く。出力はJSONのみ。"""

GEN_PROMPT_TMPL = """次のアイデアシート情報から、8枚構成のプレゼン原稿をJSONで作ってください。
情報が空の項目は、シートの他の記述から自然に補ってよいが、事実をねつ造しないこと。
各ポイントは「短い見出し」と「その説明文」をセットにすること（見出しだけで終わらせない）。

# アイデアシート情報
{idea_json}

# 出力JSON（このスキーマ厳守・```やコメント禁止）
{{
  "title": "表紙の大見出し（ロボット名を主役に、短く力強く）",
  "subtitle": "表紙のサブコピー（ひとことで魅力を伝える）",
  "team_line": "学校名・学年・氏名をまとめた1行",
  "slides": [
    {{"key": "concept", "title": "このロボットってなに？", "lead": "1行で核心",
      "points": [{{"head": "短い見出し(15字以内)", "desc": "その内容を説明する1〜2文(40〜70字)"}}]}},
    {{"key": "why",     "title": "なんで作りたいの？",   "lead": "課題を一言で", "points": [{{"head":"","desc":""}}]}},
    {{"key": "how",     "title": "どうやって動く？",     "lead": "仕組みを一言で", "points": [{{"head":"","desc":""}}]}},
    {{"key": "shape",   "title": "形と大きさ",           "lead": "見た目のポイント", "points": [{{"head":"","desc":""}}]}},
    {{"key": "who",     "title": "だれの役に立つ？",     "lead": "使う人・場面", "points": [{{"head":"","desc":""}}]}},
    {{"key": "price",   "title": "おねだんと未来",       "lead": "価格の考え方", "points": [{{"head":"","desc":""}}]}}
  ]
}}
各スライドの points は3項目（how は4項目まで可）。見出しは体言中心で簡潔に、説明文は中高生にも伝わるやさしい言葉で。"""


def generate_presentation(idea_data, model=GEN_MODEL_DEFAULT, host=None):
    """OCR結果dict → スライド原稿dict。Ollama を使用。
    コスト爆発を防ぐためサーバー側で入力文字数と出力トークンに上限を設ける。"""
    client = _ollama_client(host)
    idea_json = json.dumps(idea_data, ensure_ascii=False, indent=2)
    idea_json = idea_json[:MAX_INPUT_CHARS]   # 入力サイズの上限（巨大入力でのコスト爆発を防止）
    prompt = GEN_PROMPT_TMPL.format(idea_json=idea_json)
    resp = client.chat(
        model=model,
        messages=[
            {"role": "system", "content": GEN_SYSTEM},
            {"role": "user", "content": prompt},
        ],
        options={"temperature": 0.6, "num_predict": MAX_OUTPUT_TOKENS},  # 出力トークン上限
    )
    text = resp["message"]["content"]
    return _safe_json(text, default=_fallback_content(idea_data))


def _fallback_content(idea):
    """AI生成が失敗した場合に、OCR結果から最低限の原稿を組み立てる。"""
    name = idea.get("robot_name") or "わたしたちのロボット"

    def pts(items):
        return [{"head": str(x), "desc": ""} for x in items if str(x).strip()]

    return {
        "title": name,
        "subtitle": "アイデアをカタチに",
        "team_line": " ".join(x for x in [idea.get("school", ""), idea.get("grade", ""), idea.get("author", "")] if x),
        "slides": [
            {"key": "concept", "title": "このロボットってなに？", "lead": idea.get("impression", "")[:30],
             "points": pts([name])},
            {"key": "why", "title": "なんで作りたいの？", "lead": "", "points": pts(_split(idea.get("motivation", "")))},
            {"key": "how", "title": "どうやって動く？", "lead": "",
             "points": pts(idea.get("usage_steps", []) or _split(idea.get("shape_note", "")))},
            {"key": "shape", "title": "形と大きさ", "lead": "", "points": pts(_split(idea.get("shape_note", "")))},
            {"key": "who", "title": "だれの役に立つ？", "lead": "", "points": pts(_split(idea.get("target_demand", "")))},
            {"key": "price", "title": "おねだんと未来", "lead": "",
             "points": pts([idea.get("price", "")] if idea.get("price") else [])},
        ],
    }


# ============================================================================
# 2b. 佐藤先生方式（起承転結・各スライドの存在意義を明確化）
# ============================================================================
GEN_SYSTEM_SATO = """あなたは、東京大学名誉教授・佐藤知正先生のプレゼン指導法に従って、
高校生・高専生のロボット発表スライド（3〜5分の口頭発表）を構成する編集者です。

佐藤先生の指導原則は次の5つ:
1) 発表は「起／承／転／転続き／結準備／結」の【6フェーズ】で組み立てる。各段階に固有の「効果」がある:
   - 起 = 気を引く効果（表紙で関心をつかむ）
   - 承 = 3分間説明（用語と対象を具体的に説明する）
   - 転 = 眠気覚し（一番見せたい瞬間。動画や決定的写真で聴衆を起こす）
   - 転続き = 実体説明（仕組み・実物・動作の中身を説明）
   - 結準備 = 主張（価値を言い切る／体系化する）
   - 結 = 知識化（結論と将来課題で知識として残す）
2) 各スライドに必ず「存在意義(purpose)」(このスライドは何のためにあるのか)を一文で添える。
3) 各スライドは「視覚要素(visual_role) → それで何を説明するか(purpose)」のセットを基本単位とする。
4) 「転」では理屈を後回しにし、最もインパクトのある視覚を1枚集中させる。動画は完成後に差し替えるので、ここでは静止画＋「ここで実演します」のキャプションで作る。
5) 表紙には「概要(overview)」(この発表で何を順番に伝えるかの予告編)を入れ、聴衆の気を引く。

【文章量の指針（中高生のプレゼンに見せず、聴き応えのある分量にする）】
- overview は 100〜180字（このプレゼンで順に何を伝えるかを具体的に予告）
- 各スライドの body は 120〜200字（聴き手に語りかけるトーンの本文段落）
- 各カードの desc は 60〜120字（具体例・数字・利用シーンを入れた2〜3文）
- 与えられた【参考リサーチ】に事実・数字・類似事例があれば、自然な文中に取り込む。
  数字や事実を引用するときは文末に短く「（出典：ドメイン名）」を付ける（URL全文ではなく
  ドメイン名で簡潔に。例：「（出典：mhlw.go.jp）」）

言葉は中高生にも伝わるやさしさで、かつ発表として惹きつける流れにする。出力はJSONのみ。"""

GEN_PROMPT_SATO = """次のアイデアシート情報と参考リサーチから、佐藤先生方式（6フェーズ）の
発表スライド原稿をJSONで作ってください。
事実をねつ造せず、情報が空の項目はシートや参考リサーチから自然に補ってよい。
スライドは表紙(起)＋本文5枚＝6枚構成にすること。

# アイデアシート情報
{idea_json}

# 参考リサーチ（Google検索で得た事実・統計・類似事例。事実として使ってよいが、必ず本文に自然に溶かし込み、数字を使ったときは出典のドメイン名を短く添える）
{research}

# 出力JSON（このスキーマ厳守・```やコメント禁止）
{{
  "title": "表紙タイトル（ロボット名を主役に、短く力強く）",
  "subtitle": "表紙サブコピー（ひとことで魅力）",
  "overview": "100〜180字。このプレゼンで何を、どの順番で伝えるかを聴衆に予告する文章。リサーチに数字があればここで一つ使う",
  "team_line": "学校名・学年・氏名をまとめた1行",
  "slides": [
    {{"phase": "起", "title": "ロボット名と一番伝えたいこと",
      "purpose": "聴衆の気を引き、これから何の話をするか予告する",
      "visual_role": "表紙の魅力写真：このロボットの一番カッコいい姿",
      "lead": "ひとことキャッチコピー（20〜40字）",
      "body": "120〜200字。このロボットを作ろうと思った身近なきっかけと、その背景にある社会的な状況をリサーチの数字と組み合わせて語る",
      "points": [{{"head":"短い見出し（10〜20字）","desc":"60〜120字の説明。具体例・数字・場面を入れた2〜3文"}},
                  {{"head":"","desc":""}}]}},
    {{"phase": "承", "title": "このロボットってなに？",
      "purpose": "対象を具体的にイメージさせる（3分間説明の中核）",
      "visual_role": "ロボット全体写真または形・大きさが分かる図",
      "lead": "対象を一言で表す（20〜40字）",
      "body": "120〜200字。サイズ・見た目・特徴を具体的に。類似製品との違いがあればリサーチを参照して言及",
      "points": [{{"head":"","desc":""}},{{"head":"","desc":""}},{{"head":"","desc":""}}]}},
    {{"phase": "転", "title": "見てください、ここが見どころ",
      "purpose": "一番見せたい瞬間を出して聴衆の眠気を覚ます",
      "visual_role": "動作の決定的シーンの写真（完成後に動画に差し替える）",
      "lead": "ここで実演します",
      "body": "100〜160字。実演で何が起こるのか、どこに注目してほしいかを語りかける口調で",
      "points": [{{"head":"","desc":""}},{{"head":"","desc":""}}]}},
    {{"phase": "転続き", "title": "どうやって動くの？",
      "purpose": "仕組み・使い方の中身を実体に即して説明する",
      "visual_role": "使い方の手順図または仕組みの図解",
      "lead": "動作の流れ（20〜40字）",
      "body": "120〜200字。仕組みを段階的に説明。使われている技術（音声認識、センサー等）にリサーチで触れる",
      "points": [{{"head":"","desc":""}},{{"head":"","desc":""}},{{"head":"","desc":""}}]}},
    {{"phase": "結準備", "title": "だれの役に立つの？／なんで作りたいの？",
      "purpose": "このロボットの価値を言い切る（中心的主張）",
      "visual_role": "使われるシーンの写真または対象ユーザー",
      "lead": "価値の主張を一言で（20〜40字）",
      "body": "120〜200字。誰の、どんな困りごとを解決するか。リサーチで得た対象人口や統計を引用して規模感を示す",
      "points": [{{"head":"","desc":""}},{{"head":"","desc":""}}]}},
    {{"phase": "結", "title": "結論と、これからやりたいこと",
      "purpose": "言いたいことを言い切り、将来課題を述べて知識として残す",
      "visual_role": "なくてもよい。テキスト中心で締める",
      "lead": "まとめの一言（20〜40字）",
      "body": "120〜200字。プレゼン全体の結論と、これから改善したいこと2〜3点を熱意ある文で",
      "points": [{{"head":"","desc":""}},{{"head":"","desc":""}}]}}
  ]
}}
要件:
- 各スライドの points は 2〜3個。各 desc は 60〜120字（2〜3文）。
- 各スライドの body は 120〜200字（転は100〜160字でもよい）。
- overview は 100〜180字。
- リサーチの数字を引用したときは「（出典：ドメイン名）」を本文中に短く挿入。
- 転スライドの lead は必ず「ここで実演します」とすること。"""


def generate_presentation_sato(idea_data, model=GEN_MODEL_DEFAULT, host=None,
                                gemini_api_key=None, research_text=None):
    """佐藤先生方式（6フェーズ）の原稿を生成する。

    gemini_api_key が与えられたら、まず Gemini 2.5 Flash + Google検索で背景情報を調査し、
    その結果を Ollama 側のプロンプトに参考リサーチとして渡す。
    research_text を直接指定することもできる（呼び出し側で事前にリサーチ済みの場合）。
    """
    # 参考リサーチを取得（明示指定 > Gemini調査 > 空）
    if research_text is None and gemini_api_key:
        research_text = gemini_research_topic(idea_data, gemini_api_key)
    research_text = research_text or "（参考リサーチなし。シート情報のみから作成してください。）"

    client = _ollama_client(host)
    idea_json = json.dumps(idea_data, ensure_ascii=False, indent=2)[:MAX_INPUT_CHARS]
    # リサーチ部分は別途上限を設けて入力肥大化を防ぐ
    research_clipped = research_text[:6000]
    prompt = GEN_PROMPT_SATO.format(idea_json=idea_json, research=research_clipped)
    # 文章量を増やすので出力トークン上限も拡張
    resp = client.chat(
        model=model,
        messages=[
            {"role": "system", "content": GEN_SYSTEM_SATO},
            {"role": "user", "content": prompt},
        ],
        options={"temperature": 0.6, "num_predict": 4096},
    )
    text = resp["message"]["content"]
    content = _safe_json(text, default=_fallback_content_sato(idea_data))
    # 後段で利用できるよう調査結果も同梱（UIで確認用に出すため）
    if isinstance(content, dict):
        content.setdefault("_research", research_text)
    return content


def _fallback_content_sato(idea):
    """佐藤方式のAI生成が失敗したときの簡易フォールバック（6フェーズ・body入り）。"""
    name = idea.get("robot_name") or "わたしたちのロボット"

    def card(head, desc=""):
        return {"head": str(head).strip(), "desc": str(desc).strip()}

    def joined(items, sep="。"):
        return sep.join([s.strip() for s in items if str(s).strip()])

    motivation = (idea.get("motivation", "") or "").strip()
    impression = (idea.get("impression", "") or "").strip()
    shape = (idea.get("shape_note", "") or "").strip()
    target = (idea.get("target_demand", "") or "").strip()
    steps = [s for s in (idea.get("usage_steps", []) or []) if str(s).strip()]
    price = (idea.get("price", "") or "").strip()

    overview = (
        f"このプレゼンでは「{name}」について、なぜ作ろうと思ったのか、どんな形でどう動くのか、"
        f"そしてだれの役に立つのかを順番にお話しします。"
        + (f"{motivation[:60]}という想いが出発点です。" if motivation else "")
    )[:200]

    # 起：きっかけ
    body_kishou = (
        (impression or motivation or "")
        + "。" + (motivation if motivation and motivation != impression else "")
        + "。このロボットがあれば、毎日のちょっとした困りごとがもっと楽しく解決できると考えました。"
    ).replace("。。", "。").strip("。") + "。"

    # 承：このロボットってなに？
    body_shou_parts = []
    if shape:
        body_shou_parts.append(f"このロボットは{shape}")
    body_shou_parts.append(f"名前は「{name}」です")
    if impression:
        body_shou_parts.append(impression)
    body_shou = joined(body_shou_parts) + "。"

    # 転：見どころ
    body_ten = (
        "ここから一番見てほしい場面に入ります。"
        + (f"{steps[0]}という動きや、" if steps else "")
        + "実際に動かしてみると、思ったよりずっと表情豊かに動きます。"
        "ぜひ実演で、その様子を見てください。"
    )

    # 転続き：仕組み
    body_zoku_parts = ["このロボットは次のような手順で動きます"]
    if steps:
        body_zoku_parts.append("：" + "、".join(steps[:4]))
    body_zoku_parts.append(
        "。音声入力やセンサー、簡単なAIの仕組みを組み合わせることで、家庭にあるアラーム時計より自然なやりとりを目指しています"
    )
    body_zoku = "".join(body_zoku_parts) + "。"

    # 結準備：誰のため
    body_kj = (
        (target or "朝が苦手な人たち") + "が、毎朝もっと笑顔で起きられるようにするのが、このロボットの目的です。"
        + (f"特に{motivation[:50]}のような場面で力を発揮します。" if motivation else "")
        + "家族みんなの朝の時間を、ちょっとだけ楽しくしてくれる存在を目指しています。"
    )

    # 結
    body_ketsu = (
        f"{name}は、身近な悩みをきっかけに生まれたアイデアです。"
        + (f"想定価格は{price}を目指しています。" if price else "")
        + "これからは、もっと自然に話せる音声認識や、人によって対応を変えられる工夫を加えて、"
        "本当に使ってもらえるロボットへ育てていきたいです。"
    )

    # points（heads＋descs）を組み立てる
    p_kishou = [
        card("こんなきっかけでした",
             motivation if motivation else "毎日の困りごとを、ロボットの力で楽しく解決したいと考えました。"),
    ]
    if impression:
        p_kishou.append(card("わたしの想い", impression))

    p_shou = []
    if shape:
        p_shou.append(card("形と大きさ", shape))
    p_shou.append(card("ロボットの名前", f"このロボットは「{name}」と呼びます。覚えてもらいやすい名前を考えました。"))
    if impression and not shape:
        p_shou.append(card("印象", impression))

    p_ten = []
    if steps:
        p_ten.append(card("一番の見どころ", steps[0]))
    p_ten.append(card("注目ポイント", "動きと表情の変化に注目してください。実演でしかわからない魅力があります。"))

    p_zoku = [card(f"ステップ{i+1}", s) for i, s in enumerate(steps[:4])]
    if not p_zoku:
        p_zoku = [card("動きの仕組み", shape if shape else "センサーと音声入力を組み合わせて動きます。")]

    p_kj = []
    if target:
        # ターゲットを文として複数項目に分けず、見出しは短く保つ
        p_kj.append(card("こんな人の役に立つ", target))
    p_kj.append(card("解決したい困りごと", motivation if motivation else "毎朝のイライラした気持ちを、笑顔に変えたい。"))

    p_ketsu = []
    if price:
        p_ketsu.append(card("想定価格", f"{price}を目指しています。手の届く価格にすることで、多くの家庭で使ってもらえます。"))
    p_ketsu.append(card("これからやりたいこと",
                        "もっと自然な会話、もっと多くの表情、もっと安全な動きを実現していきたいです。"))

    return {
        "title": name,
        "subtitle": "アイデアをカタチに",
        "overview": overview,
        "team_line": " ".join(x for x in [idea.get("school", ""),
                                          idea.get("grade", ""),
                                          idea.get("author", "")] if x),
        "slides": [
            {"phase": "起", "title": name,
             "purpose": "聴衆の気を引き、これから何の話をするか予告する",
             "visual_role": "表紙の魅力写真",
             "lead": "身近な悩みから生まれたロボットです",
             "body": body_kishou, "points": p_kishou[:3]},
            {"phase": "承", "title": "このロボットってなに？",
             "purpose": "対象を具体的にイメージさせる（3分間説明）",
             "visual_role": "ロボット全体の写真または形・大きさが分かる図",
             "lead": "形・大きさ・特徴を紹介します",
             "body": body_shou, "points": p_shou[:3]},
            {"phase": "転", "title": "ここが見どころ",
             "purpose": "一番見せたい瞬間を出して聴衆の眠気を覚ます",
             "visual_role": "動作の決定的シーンの写真（完成後に動画に差し替え）",
             "lead": "ここで実演します",
             "body": body_ten, "points": p_ten[:2]},
            {"phase": "転続き", "title": "どうやって動くの？",
             "purpose": "仕組み・使い方の中身を実体に即して説明する",
             "visual_role": "使い方の手順図または仕組みの図解",
             "lead": "動作の流れを順番に",
             "body": body_zoku, "points": p_zoku[:3]},
            {"phase": "結準備", "title": "だれの役に立つの？",
             "purpose": "このロボットの価値を言い切る（中心的主張）",
             "visual_role": "使われるシーンの写真または対象ユーザー",
             "lead": "こんな人の毎日を変えたい",
             "body": body_kj, "points": p_kj[:3]},
            {"phase": "結", "title": "結論と、これからやりたいこと",
             "purpose": "言いたいことを言い切り、将来課題を述べて知識として残す",
             "visual_role": "テキスト中心で締める",
             "lead": "まとめと未来への一歩",
             "body": body_ketsu, "points": p_ketsu[:2]},
        ],
    }


# ============================================================================
# 2.5) 原稿生成（business調・提案書スタイル）
# ============================================================================
# 7セクション構成：表紙 + Problem → Solution → How → Impact → Why → Next
BIZ_SECTIONS = ["problem", "solution", "how", "impact", "why", "next"]
BIZ_SECTION_LABELS = {
    "problem":  ("Problem & Market Context", "課題と市場背景"),
    "solution": ("Solution Overview",        "ソリューション概要"),
    "how":      ("How It Works",             "仕組みと差別化"),
    "impact":   ("Use Cases & Impact",       "活用シーンと効果"),
    "why":      ("Why Now / Why Us",         "今、なぜ我々か"),
    "next":     ("Next Steps",               "次のアクション"),
}

GEN_SYSTEM_BIZ = """あなたは、ロボット技術の事業提案書を作成するシニアコンサルタントです。
聴き手は審査員・投資家・技術者など、課題解決の規模と効果に関心がある大人。

【書き方の指針】
- 中高生向けのやさしい口調ではなく、ビジネス文書の冷静・簡潔・断定的なトーン
- 「〜したい」より「〜する」「〜できる」と言い切る
- 数値・統計・proof point（具体的事実）を積極的に使う
- 解決すべき課題の規模感を最初に提示する（Problem → Solution → Impact）
- 専門用語は使ってよいが、必ず文脈で意味が分かるように
- 「気持ち」「想い」より「効果」「価値」「市場機会」を語る

【セクション構成】(必ずこの順序で7枚=表紙+6枚)
1. 表紙 — 提案タイトル＋サブコピー
2. Problem & Market Context — 解決すべき課題と市場規模/統計
3. Solution Overview — このロボットが何を解決するか
4. How It Works — 仕組み・技術的差別化（SLAM、AI、センサー等の用語OK）
5. Use Cases & Impact — 具体的利用シーン3つ＋効果（KPI数値）
6. Why Now / Why Us — タイミング・競合優位・実現可能性
7. Next Steps — 提案する次の検証ステップ

【KPIの扱い】
各スライドに、可能なら kpis フィールドで「大きな数字とラベル」を入れる:
  例: kpis: [{"value": "48歳", "label": "ドライバー平均年齢", "source": "mlit.go.jp"}]
出典はドメイン名のみ短く。リサーチに数字があれば必ず使う。

【文章量】
- overview: 100〜180字
- 各スライド body: 130〜220字
- 各カード desc: 60〜120字

出力はJSONのみ。"""

GEN_PROMPT_BIZ = """次のアイデアシートと参考リサーチから、business調（提案書スタイル）の
プレゼン原稿をJSONで作ってください。
表紙＋本文6枚＝合計7枚構成、セクションは problem→solution→how→impact→why→next の順に固定。

# アイデアシート情報
{idea_json}

# 参考リサーチ（Google検索で得た事実・統計・類似事例。本文中に「（出典：ドメイン名）」を短く付ける）
{research}

# 出力JSON（このスキーマ厳守・```やコメント禁止）
{{
  "title": "提案タイトル（ロボット名を含む、価値訴求型。例：『荷役負担をゼロにするアーム搭載型トラック』）",
  "subtitle": "サブコピー（誰の何を解決するかを一行で）",
  "overview": "100〜180字。この提案で訴えるストーリー全体を予告。市場規模の数字を1つ含めると良い",
  "team_line": "学校名・学年・氏名をまとめた1行",
  "slides": [
    {{"section": "problem",
      "title": "課題と市場背景",
      "purpose": "解決すべき課題の規模と緊急性を示す",
      "lead": "見出しコピー（25〜45字。課題を端的に）",
      "body": "130〜220字。課題の現状と背景、なぜいま解決すべきかを論理的に述べる。リサーチの統計を必ず引用",
      "kpis": [
        {{"value": "数字", "label": "ラベル", "source": "ドメイン"}},
        {{"value": "数字", "label": "ラベル", "source": "ドメイン"}}
      ],
      "points": [{{"head":"短い小見出し","desc":"60〜120字の具体説明"}}]
    }},
    {{"section": "solution",
      "title": "ソリューション概要",
      "purpose": "このロボットが課題をどう解決するかを一文で言い切る",
      "lead": "価値提案を一言で",
      "body": "130〜220字。何ができるロボットか、誰の何が変わるかを断定的に",
      "kpis": [],
      "points": [{{"head":"提供価値1","desc":""}},
                  {{"head":"提供価値2","desc":""}}]
    }},
    {{"section": "how",
      "title": "仕組みと差別化",
      "purpose": "技術的にどう実現するか、なぜ既存品より優れるかを示す",
      "lead": "技術アプローチを一言で",
      "body": "130〜220字。要素技術や設計思想、既存ソリューションとの違い",
      "kpis": [],
      "points": [{{"head":"技術1","desc":""}},
                  {{"head":"技術2","desc":""}},
                  {{"head":"技術3","desc":""}}]
    }},
    {{"section": "impact",
      "title": "活用シーンと効果",
      "purpose": "誰の業務がどう変わるか、効果を数値で示す",
      "lead": "効果を見出しで",
      "body": "130〜220字。具体的なユースケース描写と、もたらす変化",
      "kpis": [
        {{"value": "数字", "label": "改善指標", "source": ""}},
        {{"value": "数字", "label": "改善指標", "source": ""}}
      ],
      "points": [{{"head":"ユースケース1","desc":""}},
                  {{"head":"ユースケース2","desc":""}}]
    }},
    {{"section": "why",
      "title": "今、なぜ我々か",
      "purpose": "市場タイミングと、自分たちが解く優位性",
      "lead": "Why now / why us を一言で",
      "body": "130〜220字。市場の追い風、競合との違い、実現可能性",
      "kpis": [],
      "points": [{{"head":"差別化1","desc":""}},
                  {{"head":"差別化2","desc":""}}]
    }},
    {{"section": "next",
      "title": "次のアクション",
      "purpose": "聴き手に何をしてほしいかを明示する",
      "lead": "提案する次のステップ",
      "body": "130〜220字。検証フェーズの計画、必要な協力、想定タイムライン",
      "kpis": [
        {{"value": "想定価格", "label": "コスト", "source": ""}}
      ],
      "points": [{{"head":"次のステップ1","desc":""}},
                  {{"head":"次のステップ2","desc":""}}],
      "cta": "具体的なお願い・依頼を一言で"
    }}
  ]
}}

要件:
- セクション順は厳守
- 中高生口調・絵文字・「すごい」「面白い」などの主観表現は避ける
- 数値があれば必ずKPIに入れる
- proof pointの出典は「（出典：ドメイン名）」形式で本文中に
"""


def generate_presentation_biz(idea_data, model=GEN_MODEL_DEFAULT, host=None,
                               gemini_api_key=None, research_text=None):
    """business調（提案書スタイル）の原稿を生成する。

    gemini_api_key が指定されたら Web リサーチで背景情報を補強。
    research_text を直接渡すこともできる。
    """
    # 参考リサーチを取得（明示指定 > Gemini調査 > 空）
    if research_text is None and gemini_api_key:
        research_text = gemini_research_topic(idea_data, gemini_api_key)
    research_text = research_text or "（参考リサーチなし。シート情報のみで作成してください。）"

    client = _ollama_client(host)
    idea_json = json.dumps(idea_data, ensure_ascii=False, indent=2)[:MAX_INPUT_CHARS]
    research_clipped = research_text[:6000]
    prompt = GEN_PROMPT_BIZ.format(idea_json=idea_json, research=research_clipped)
    resp = client.chat(
        model=model,
        messages=[
            {"role": "system", "content": GEN_SYSTEM_BIZ},
            {"role": "user", "content": prompt},
        ],
        options={"temperature": 0.4, "num_predict": 4096},
    )
    text = resp["message"]["content"]
    content = _safe_json(text, default=_fallback_content_biz(idea_data))
    if isinstance(content, dict):
        content.setdefault("_research", research_text)
    return content


def _fallback_content_biz(idea):
    """business方式のAI生成失敗時フォールバック。7枚構成のbody入りリッチ版。"""
    name = idea.get("robot_name") or "本ロボット"

    def card(head, desc=""):
        return {"head": str(head).strip(), "desc": str(desc).strip()}

    motivation = (idea.get("motivation", "") or "").strip()
    impression = (idea.get("impression", "") or "").strip()
    shape = (idea.get("shape_note", "") or "").strip()
    target = (idea.get("target_demand", "") or "").strip()
    steps = [s for s in (idea.get("usage_steps", []) or []) if str(s).strip()]
    price = (idea.get("price", "") or "").strip()

    overview = (
        f"本提案では「{name}」を用いた課題解決のアプローチを示す。"
        f"{(motivation or '')[:80]}という現状に対し、" if motivation else f"対象領域の現状に対し、"
    )[:200] + "本ロボットがもたらす効果と、次の検証ステップを順に述べる。"

    return {
        "title": f"{name}による課題解決の提案",
        "subtitle": ((target[:40] + "向けの新しい選択肢") if target else "新しいロボット技術による業務改善"),
        "overview": overview,
        "team_line": " ".join(x for x in [idea.get("school", ""),
                                          idea.get("grade", ""),
                                          idea.get("author", "")] if x),
        "slides": [
            {"section": "problem",
             "title": "課題と市場背景",
             "purpose": "解決すべき課題の規模と緊急性を示す",
             "lead": (motivation[:40] if motivation else "現状の課題と背景"),
             "body": (motivation or "現状の業務には改善余地がある。") +
                     "本提案は、この課題を技術的に解決する手段として本ロボットを位置づける。",
             "kpis": [],
             "points": [card("現状の課題", motivation or "業務負担の増大が継続的な問題となっている。"),
                        card("背景", impression or "現場でのヒアリングから課題の深刻さが確認できた。")][:3]},
            {"section": "solution",
             "title": "ソリューション概要",
             "purpose": "このロボットが何を解決するかを言い切る",
             "lead": f"{name}が課題を解決する",
             "body": f"{name}は、" + (shape if shape else "対象作業に最適化された設計") +
                     "を持ち、対象者の作業負担を直接的に軽減する。",
             "kpis": [],
             "points": [card("提供価値", "対象作業の自動化または半自動化"),
                        card("適用範囲", target or "対象ユーザー全般")][:3]},
            {"section": "how",
             "title": "仕組みと差別化",
             "purpose": "技術的実現と既存品との差別化を示す",
             "lead": "技術アプローチ",
             "body": ("動作フロー：" + "／".join(steps[:4]) + "。" if steps else "") +
                     "センサーと制御の連携により、既存ソリューションでは難しかった範囲をカバーする。",
             "kpis": [],
             "points": [card(f"ステップ{i+1}", s) for i, s in enumerate(steps[:3])] or
                       [card("設計思想", shape or "対象作業に最適化")]
             },
            {"section": "impact",
             "title": "活用シーンと効果",
             "purpose": "業務がどう変わるかを具体化",
             "lead": "対象者の業務を変える",
             "body": f"{target or '対象ユーザー'}の作業環境において、本ロボット導入後は作業時間の短縮と負担軽減が見込まれる。",
             "kpis": ([{"value": price, "label": "想定価格", "source": ""}] if price else []),
             "points": [card("ユースケース1", target or "対象業務での導入"),
                        card("ユースケース2", "派生的な活用範囲の拡張")][:3]},
            {"section": "why",
             "title": "今、なぜ我々か",
             "purpose": "タイミングと優位性",
             "lead": "Why now / Why us",
             "body": "ロボット技術の普及と現場の負担増の両面から、本提案は今このタイミングでの導入価値が高い。" +
                     "本チームは現場視点と技術視点の両方を持つ。",
             "kpis": [],
             "points": [card("市場タイミング", "技術成熟と現場ニーズの一致"),
                        card("優位性", "現場視点を持つ設計")][:3]},
            {"section": "next",
             "title": "次のアクション",
             "purpose": "聴き手に何をしてほしいかを明示",
             "lead": "次の検証ステップ",
             "body": "試作機の評価、現場での実証、対象ユーザーからのフィードバック収集を順に進める提案である。",
             "kpis": ([{"value": price, "label": "想定単価", "source": ""}] if price else []),
             "points": [card("検証ステップ1", "試作機の評価"),
                        card("検証ステップ2", "現場ユーザーへのヒアリング")][:3],
             "cta": "実証フェーズへの協力を依頼"},
        ],
    }


# ============================================================================
# 3. PowerPoint 生成（python-pptx）
# ============================================================================
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
# フチに触れさせない安全マージン（複合機が端まで印刷できないため）
SAFE = Inches(0.55)


def _rgb(hexstr):
    return RGBColor.from_string(hexstr)


def _set_white_bg(slide):
    """スライド背景を必ず白に。"""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb("FFFFFF")


def _set_bg(slide, hexcolor):
    """スライド背景を任意色に設定する。"""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(hexcolor)


def _no_line(shape):
    shape.line.fill.background()


def _txt(slide, l, t, w, h, lines, *, size, color, bold=False, font=FONT_BODY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, leading=1.12, space_after=4):
    """テキストボックス追加。lines は str か (str, opts) のリスト。"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        opts = {}
        if isinstance(ln, tuple):
            ln, opts = ln
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.line_spacing = leading
        p.space_after = Pt(opts.get("space_after", space_after))
        p.space_before = Pt(0)
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", bold)
        r.font.name = opts.get("font", font)
        r.font.color.rgb = _rgb(opts.get("color", color))
        _set_cjk_font(r, opts.get("font", font))
    return tb


def _set_cjk_font(run, font_name):
    """日本語が確実にそのフォントで出るよう east-asian 指定を入れる。"""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font_name)


def _badge(slide, l, t, d, text, color):
    """連番などを入れる色付き円バッジ。"""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, d, d)
    c.fill.solid()
    c.fill.fore_color.rgb = _rgb(color)
    _no_line(c)
    tf = c.text_frame
    tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(16); r.font.bold = True; r.font.name = FONT_HEAD
    r.font.color.rgb = _rgb("FFFFFF")
    _set_cjk_font(r, FONT_HEAD)
    return c


def _card(slide, l, t, w, h, fill=COL_BG_CARD):
    """角丸の薄いカード（中に置くだけ。端には触れさせない）。"""
    rr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    rr.fill.solid()
    rr.fill.fore_color.rgb = _rgb(fill)
    _no_line(rr)
    try:
        rr.adjustments[0] = 0.06
    except Exception:
        pass
    return rr


def _title(slide, text, accent):
    """各スライド共通のタイトル（アクセント色のドット＋見出し）。下線は引かない。"""
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, SAFE, Inches(0.62), Inches(0.20), Inches(0.20))
    dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(accent); _no_line(dot)
    _txt(slide, SAFE + Inches(0.34), Inches(0.5), SLIDE_W - SAFE * 2 - Inches(0.34), Inches(0.7),
         text, size=30, color=COL_INK, bold=True, font=FONT_HEAD)


def _add_picture_fit(slide, img_bytes, l, t, max_w, max_h):
    """画像を枠内にアスペクト比維持で収め、中央寄せで配置。"""
    from PIL import Image
    bio = io.BytesIO(img_bytes)
    iw, ih = Image.open(bio).size
    ratio = min(max_w / iw, max_h / ih)
    w = int(iw * ratio); h = int(ih * ratio)
    l2 = l + (max_w - w) // 2
    t2 = t + (max_h - h) // 2
    slide.shapes.add_picture(io.BytesIO(img_bytes), l2, t2, width=w, height=h)


# 薄いアクセント色（白に混ぜたティント）
_TINTS = {COL_A1: "FCEDE6", COL_A2: "E7EEFD", COL_A3: "E6F4EC"}


def _norm_point(p):
    """points要素を (見出し, 説明) に正規化。文字列・dictどちらも許容。"""
    if isinstance(p, dict):
        return (str(p.get("head") or p.get("text") or "").strip(),
                str(p.get("desc") or "").strip())
    return (str(p).strip(), "")


def _point_card(slide, l, t, w, h, n, head, desc, accent):
    """番号バッジ＋見出し＋説明文の1枚カード（薄い背景・角丸）。"""
    _card(slide, l, t, w, h, fill="F5F6F8")
    bd = Inches(0.5)
    bx = l + Inches(0.26)
    if desc:
        by = t + Inches(0.24)
        head_anchor = MSO_ANCHOR.TOP
    else:
        by = t + (h - bd) // 2          # 説明が無い時は縦中央
        head_anchor = MSO_ANCHOR.MIDDLE
    _badge(slide, bx, by, bd, str(n), accent)
    tx = bx + bd + Inches(0.28)
    tw = l + w - tx - Inches(0.3)
    if desc:
        _txt(slide, tx, t + Inches(0.2), tw, Inches(0.42),
             head, size=17, color=COL_INK, bold=True, font=FONT_HEAD, leading=1.05)
        _txt(slide, tx, t + Inches(0.66), tw, h - Inches(0.78),
             desc, size=12.5, color=COL_MUTE, font=FONT_BODY, leading=1.2)
    else:
        _txt(slide, tx, t, tw, h, head, size=17, color=COL_INK, bold=True,
             font=FONT_HEAD, anchor=MSO_ANCHOR.MIDDLE, leading=1.1)


def _cards_column(slide, points, l, t, w, accent, bottom=Inches(7.0), max_n=4):
    """縦方向にカードを敷き詰める。"""
    pts = [_norm_point(p) for p in points if _norm_point(p)[0]][:max_n]
    if not pts:
        return
    n = len(pts)
    gap = Inches(0.22)
    avail = bottom - t
    ch = int((avail - gap * (n - 1)) / n)
    ch = max(min(ch, Inches(1.55)), Inches(0.9))
    for i, (head, desc) in enumerate(pts):
        cy = t + (ch + gap) * i
        _point_card(slide, l, cy, w, ch, i + 1, head, desc, accent)


def _render_split(slide, points, l_text_w, body_top, img_bytes, accent, max_n=3):
    """左テキスト / 右画像 のスプリット描画（どの本文スライドでも使える）。"""
    _cards_column(slide, points, SAFE, body_top, l_text_w, accent, max_n=max_n)
    pic_l = SAFE + l_text_w + Inches(0.35)
    pic_w = SLIDE_W - pic_l - SAFE
    pic_h = Inches(7.0) - body_top
    _card(slide, pic_l, body_top, pic_w, pic_h, fill=COL_BG_CARD)
    _add_picture_fit(slide, img_bytes, pic_l + Inches(0.2), body_top + Inches(0.2),
                     pic_w - Inches(0.4), pic_h - Inches(0.4))


# 画像の自動振り分け優先順（写真リストだけ渡された場合）
_AUTO_KEYS = ["shape", "concept", "how", "who", "why"]


def _photos_to_map(photos):
    m = {}
    for img, key in zip(photos, _AUTO_KEYS):
        m[key] = img
    return m


def build_pptx(content, photos=None, output_path="robot_presentation.pptx",
               price_value="", images_by_key=None):
    """
    content       : generate_presentation() の戻り dict
    images_by_key : {本文スライドkey: 画像bytes} の対応（優先）。
    photos        : 画像bytesのリスト（後方互換。指定優先順で各ページへ自動配分）。
    ※ ギャラリーページは作らない。割り当てられた画像は各ページにスプリット配置。
    """
    if images_by_key is None:
        images_by_key = _photos_to_map(photos or [])
    images_by_key = {k: v for k, v in images_by_key.items() if v}
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ---- 表紙（必ず白背景・ビジネス調を避ける） ----
    s = prs.slides.add_slide(blank)
    _set_white_bg(s)
    # 上部にカラフルな丸を3つ（端には触れない位置に）— 競技会らしい遊び心
    for i, col in enumerate(ACCENTS):
        d = Inches(0.5)
        c = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               SAFE + Inches(0.0) + Inches(0.7) * i, Inches(1.05), d, d)
        c.fill.solid(); c.fill.fore_color.rgb = _rgb(col); _no_line(c)
    # カテゴリタグ（角丸・中置き）
    tag = _card(s, SAFE, Inches(1.9), Inches(3.4), Inches(0.5), fill=COL_INK)
    _txt(s, SAFE, Inches(1.97), Inches(3.4), Inches(0.36),
         "ロボット甲子園 アイデア提案", size=14, color="FFFFFF", bold=True,
         font=FONT_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 大見出し
    _txt(s, SAFE, Inches(2.75), SLIDE_W - SAFE * 2, Inches(2.0),
         content.get("title", ""), size=54, color=COL_INK, bold=True, font=FONT_HEAD)
    # サブコピー（アクセント色）
    _txt(s, SAFE, Inches(4.7), SLIDE_W - SAFE * 2, Inches(0.8),
         content.get("subtitle", ""), size=22, color=COL_A1, bold=True, font=FONT_HEAD)
    # チーム情報
    _txt(s, SAFE, Inches(6.35), SLIDE_W - SAFE * 2, Inches(0.6),
         content.get("team_line", ""), size=16, color=COL_MUTE, font=FONT_BODY)

    # ---- 本文スライド ----
    slides = content.get("slides", [])
    for idx, sl in enumerate(slides):
        s = prs.slides.add_slide(blank)
        _set_white_bg(s)
        accent = ACCENTS[idx % len(ACCENTS)]
        _title(s, sl.get("title", ""), accent)

        key = sl.get("key", "")
        lead = sl.get("lead", "")
        points = sl.get("points") or []
        img = images_by_key.get(key)

        # リード文（共通）
        if lead:
            _txt(s, SAFE, Inches(1.55), SLIDE_W - SAFE * 2, Inches(0.6),
                 lead, size=19, color=accent, bold=True, font=FONT_HEAD)
            body_top = Inches(2.35)
        else:
            body_top = Inches(1.95)

        # ---- レイアウト分岐 ----
        if key == "price":
            # 左に金額パネル / 右に補足カード（価格ページは画像より金額を主役に）
            panel_w = Inches(4.55)
            panel_h = Inches(7.0) - body_top
            price_txt = price_value or _first_price(points)
            _card(s, SAFE, body_top, panel_w, panel_h, fill=_TINTS.get(accent, COL_BG_CARD))
            _txt(s, SAFE + Inches(0.4), body_top + Inches(0.5), panel_w - Inches(0.8), Inches(0.5),
                 "想定価格", size=15, color=accent, bold=True, font=FONT_HEAD)
            _txt(s, SAFE + Inches(0.4), body_top + Inches(1.15), panel_w - Inches(0.8), panel_h - Inches(1.5),
                 price_txt, size=30, color=COL_INK, bold=True, font=FONT_HEAD, leading=1.15)
            _cards_column(s, points, SAFE + panel_w + Inches(0.55), body_top,
                          SLIDE_W - (SAFE + panel_w + Inches(0.55)) - SAFE, accent, max_n=3)

        elif img:
            # 画像が割り当てられたページ：左テキスト / 右画像
            _render_split(s, points, Inches(5.85), body_top, img, accent, max_n=3)

        else:
            # 全幅カード（見出し＋説明）
            max_n = 4 if key == "how" else 3
            _cards_column(s, points, SAFE, body_top, SLIDE_W - SAFE * 2, accent, max_n=max_n)

    # ---- 裏表紙（白・シンプル） ----
    s = prs.slides.add_slide(blank)
    _set_white_bg(s)
    _txt(s, SAFE, Inches(3.0), SLIDE_W - SAFE * 2, Inches(1.0),
         "ありがとうございました！", size=40, color=COL_INK, bold=True,
         font=FONT_HEAD, align=PP_ALIGN.CENTER)
    _txt(s, SAFE, Inches(4.1), SLIDE_W - SAFE * 2, Inches(0.6),
         content.get("team_line", ""), size=16, color=COL_MUTE,
         font=FONT_BODY, align=PP_ALIGN.CENTER)

    prs.save(output_path)
    return output_path


# 佐藤先生方式：6フェーズ＋効果ラベル＋配色
_PHASE_ORDER = ["起", "承", "転", "転続き", "結準備", "結"]
_PHASE_EFFECTS = {
    "起":   "気を引く効果",
    "承":   "3分間説明",
    "転":   "眠気覚し",
    "転続き": "実体説明",
    "結準備": "主張",
    "結":   "知識化",
}
_PHASE_COLORS = {
    "起":   COL_A2,    # ブルー（落ち着きと期待）
    "承":   "0891B2",  # シアン（説明）
    "転":   "DC2626",  # レッド（眠気覚し・最重要）
    "転続き": "EA580C",  # オレンジ（勢いを維持）
    "結準備": COL_A3,    # グリーン（主張の確立）
    "結":   "7C3AED",  # パープル（知識化）
}
_PHASE_TINTS = {
    "起":   "E7EEFD",
    "承":   "E0F7FA",
    "転":   "FEE2E2",
    "転続き": "FFEDD5",
    "結準備": "E6F4EC",
    "結":   "EDE9FE",
}

# 旧4幕用（互換のため残す。新しい実装からは参照しない）
_ACT_COLORS = {"起": COL_A1, "承": COL_A2, "転": COL_A3, "結": COL_INK}


def _phase_of(slide_dict):
    """新スキーマ(phase)優先、旧スキーマ(act)もフォールバック許容。"""
    return (slide_dict.get("phase") or slide_dict.get("act") or "").strip()


def build_pptx_sato(content, photos=None, output_path="robot_presentation_sato.pptx",
                    images_by_key=None, cover_image=None):
    """佐藤先生方式（6フェーズ・各スライドの存在意義つき）のPPTXを生成する。

    フェーズ: 起 → 承 → 転 → 転続き → 結準備 → 結（合計6枚＋表紙裏表紙）
    images_by_key: {スライドindex: 画像bytes}。
                   未指定なら photos から 転 → 転続き → 起 → 承 → 結準備 の順で自動配分。
    cover_image:   表紙に置く画像/アイコンの bytes（任意）。指定すると右側にヒーロー画像として配置。
                   未指定なら3色ドットのデフォルト装飾。
    各スライドは「body段落＋カード」の2層構造で文章量を確保。
    転スライドには「🎬 ここで実演します」のキャプションを必ず入れる。
    """
    slides = content.get("slides", [])

    # 画像の割り当て：明示が無ければ「転」「転続き」を最優先
    if images_by_key is None:
        images_by_key = {}
        photos = photos or []
        if photos:
            priority_phase_order = ["転", "転続き", "起", "承", "結準備", "結"]
            ordered_indices = []
            for ph in priority_phase_order:
                for i, sl in enumerate(slides):
                    if _phase_of(sl) == ph and i not in ordered_indices:
                        ordered_indices.append(i)
            for img, i in zip(photos, ordered_indices):
                images_by_key[i] = img
    images_by_key = {k: v for k, v in images_by_key.items() if v}

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # =====================================================================
    # 表紙
    # =====================================================================
    s = prs.slides.add_slide(blank)
    _set_white_bg(s)

    # 表紙アイコン/写真がアップされていれば右半分にヒーロー画像
    if cover_image:
        hero_l = SLIDE_W - SAFE - Inches(4.8)
        hero_t = Inches(1.0)
        hero_w = Inches(4.8)
        hero_h = Inches(5.5)
        _card(s, hero_l, hero_t, hero_w, hero_h, fill=COL_BG_CARD)
        try:
            _add_picture_fit(s, cover_image, hero_l + Inches(0.2), hero_t + Inches(0.2),
                             hero_w - Inches(0.4), hero_h - Inches(0.4))
        except Exception:
            pass  # 画像読み込みエラー時はカードだけ残す
        text_right = hero_l - Inches(0.3)
    else:
        # 3色ドットのデフォルト装飾
        for i, col in enumerate([_PHASE_COLORS["起"], _PHASE_COLORS["転"], _PHASE_COLORS["結準備"]]):
            d = Inches(0.5)
            c = s.shapes.add_shape(MSO_SHAPE.OVAL, SAFE + Inches(0.7) * i, Inches(1.0), d, d)
            c.fill.solid(); c.fill.fore_color.rgb = _rgb(col); _no_line(c)
        text_right = SLIDE_W - SAFE

    # 「ロボット甲子園 発表スライド」タグ
    _card(s, SAFE, Inches(1.85), Inches(3.4), Inches(0.5), fill=COL_INK)
    _txt(s, SAFE, Inches(1.92), Inches(3.4), Inches(0.36),
         "ロボット甲子園 発表スライド", size=14, color="FFFFFF", bold=True,
         font=FONT_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # メインタイトル（画像があれば左半分のみに収める）
    title_w = text_right - SAFE
    _txt(s, SAFE, Inches(2.6), title_w, Inches(1.4),
         content.get("title", ""), size=46, color=COL_INK, bold=True, font=FONT_HEAD)
    # サブタイトル
    _txt(s, SAFE, Inches(4.2), title_w, Inches(0.7),
         content.get("subtitle", ""), size=20, color=COL_A1, bold=True, font=FONT_HEAD)
    # 概要パネル（佐藤先生方式：talkの予告編）
    ov = content.get("overview", "")
    if ov:
        ov_w = text_right - SAFE
        _card(s, SAFE, Inches(5.0), ov_w, Inches(1.5), fill=COL_BG_CARD)
        _txt(s, SAFE + Inches(0.3), Inches(5.12), ov_w - Inches(0.6), Inches(0.4),
             "📋 概要（このプレゼンで伝えること）",
             size=13, color=COL_MUTE, bold=True, font=FONT_HEAD)
        _txt(s, SAFE + Inches(0.3), Inches(5.5), ov_w - Inches(0.6), Inches(0.95),
             ov, size=13, color=COL_INK, font=FONT_BODY, leading=1.35)
    # チームライン
    _txt(s, SAFE, Inches(6.75), SLIDE_W - SAFE * 2, Inches(0.45),
         content.get("team_line", ""), size=14, color=COL_MUTE, font=FONT_BODY)

    # =====================================================================
    # 本文（6フェーズ）
    # =====================================================================
    # フェーズバッジの幅マップ（「結準備」「転続き」は3文字なので広め）
    _BADGE_W = {"起": Inches(0.7), "承": Inches(0.7), "転": Inches(0.7),
                "転続き": Inches(1.05), "結準備": Inches(1.05), "結": Inches(0.7)}
    _BADGE_FONT = {"起": 24, "承": 24, "転": 24,
                   "転続き": 14, "結準備": 14, "結": 24}

    # 日本語行高さ（pt → インチ）: 13pt × leading 1.45 ≒ 0.262 in
    LINE_H_BODY_IN = 13 * 1.45 / 72
    LINE_H_DESC_IN = 12.5 * 1.25 / 72
    # 1行に入る日本語文字数（CJKは1文字≒13pt幅と概算、半角換算は0.5扱い）
    def _chars_per_line(width_emu, font_pt=13):
        in_w = width_emu / 914400.0  # EMU→インチ
        return max(10, int(in_w * 72 / font_pt))

    def _est_lines(text, width_emu, font_pt=13):
        if not text:
            return 0
        # 改行は別カウント、長い文は折り返し計算
        cpl = _chars_per_line(width_emu, font_pt)
        lines = 0
        for para in text.split("\n"):
            n = max(1, (len(para) + cpl - 1) // cpl)
            lines += n
        return max(1, lines)

    for idx, sl in enumerate(slides):
        s = prs.slides.add_slide(blank)
        _set_white_bg(s)
        phase = _phase_of(sl)
        accent = _PHASE_COLORS.get(phase, COL_A2)
        tint = _PHASE_TINTS.get(phase, COL_BG_CARD)
        effect = _PHASE_EFFECTS.get(phase, "")

        # ─── 上部：フェーズバッジ＋効果ラベル＋タイトル ───
        badge_w = _BADGE_W.get(phase, Inches(0.7))
        badge_font = _BADGE_FONT.get(phase, 20)
        if phase:
            _card(s, SAFE, Inches(0.5), badge_w, Inches(0.7), fill=accent)
            _txt(s, SAFE, Inches(0.5), badge_w, Inches(0.7),
                 phase, size=badge_font, color="FFFFFF", bold=True, font=FONT_HEAD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if effect:
                _txt(s, SAFE + badge_w + Inches(0.15), Inches(0.5), Inches(2.2), Inches(0.32),
                     effect, size=11, color=accent, bold=True, font=FONT_HEAD)
            title_l = SAFE + badge_w + Inches(0.15)
            title_t = Inches(0.85)
        else:
            title_l = SAFE
            title_t = Inches(0.5)

        _txt(s, title_l, title_t, SLIDE_W - title_l - SAFE, Inches(0.55),
             sl.get("title", ""), size=24, color=COL_INK, bold=True, font=FONT_HEAD)

        # ─── 存在意義バー ───
        purpose = sl.get("purpose", "")
        top = Inches(1.55)
        if purpose:
            _card(s, SAFE, top, SLIDE_W - SAFE * 2, Inches(0.5), fill=tint)
            _txt(s, SAFE + Inches(0.25), top, SLIDE_W - SAFE * 2 - Inches(0.5), Inches(0.5),
                 "🎯 このスライドの役割：" + purpose,
                 size=12, color=COL_INK, bold=True, font=FONT_HEAD,
                 anchor=MSO_ANCHOR.MIDDLE)
            top = top + Inches(0.6)

        # ─── レイアウト寸法 ───
        points = sl.get("points") or []
        img = images_by_key.get(idx)
        visual_role = sl.get("visual_role", "")
        body_text = sl.get("body", "")
        lead = sl.get("lead", "")
        if phase == "転" and not lead:
            lead = "ここで実演します"

        # テキスト領域の幅（画像があれば左半分、なければ全幅）
        if img or phase == "転":
            text_w = Inches(5.9)
            img_l = SAFE + text_w + Inches(0.35)
            img_w = SLIDE_W - img_l - SAFE
        else:
            text_w = SLIDE_W - SAFE * 2
            img_l = None
            img_w = None

        cards_bottom = Inches(7.15)  # スライド下端 7.5 から余白 0.35
        available_h = cards_bottom - top  # この縦範囲に lead/body/cards 全部が収まる必要がある

        # ─── 各要素の必要高さを見積もる ───
        lead_h = Inches(0.45) if lead else 0
        body_lines = _est_lines(body_text, text_w, font_pt=13) if body_text else 0
        # body 高さは行数 × 行高 ＋ 余白
        body_h_needed = (Inches(LINE_H_BODY_IN * body_lines + 0.10)
                         if body_text else 0)
        # body の上限：使える縦の35%まで（カード領域を確保）
        body_h_max = int(available_h * 0.35)
        if body_h_needed and body_h_needed > body_h_max:
            body_h_needed = body_h_max  # 上限でクリップ（はみ出し防止）

        # ─── 描画：リード ───
        if lead:
            _txt(s, SAFE, top, text_w, lead_h,
                 lead, size=16, color=accent, bold=True, font=FONT_HEAD)
            top = top + lead_h + Inches(0.08)

        # ─── 描画：body段落 ───
        if body_text:
            _txt(s, SAFE, top, text_w, body_h_needed,
                 body_text, size=13, color=COL_INK, font=FONT_BODY, leading=1.4)
            top = top + body_h_needed + Inches(0.18)

        # ─── 描画：ポイントカード（descの文字数に応じて高さを動的計算） ───
        if points and top < cards_bottom - Inches(0.4):
            # まず、各カードに必要な高さを desc から計算
            pts = []
            for p in points[:3]:
                if isinstance(p, dict):
                    head = (p.get("head") or "").strip()
                    desc = (p.get("desc") or "").strip()
                else:
                    head, desc = str(p).strip(), ""
                if not head and not desc:
                    continue
                # desc用の幅は左に番号バッジ(0.5)＋余白(0.84)を引いた残り
                desc_w = text_w - Inches(1.34)
                if desc:
                    desc_lines = _est_lines(desc, desc_w, font_pt=12.5)
                    # 見出し1行(0.42) + 余白(0.22) + desc + 下マージン(0.18)
                    ch_needed = Inches(0.42 + 0.22 + LINE_H_DESC_IN * desc_lines + 0.20)
                else:
                    ch_needed = Inches(0.65)
                pts.append((head, desc, ch_needed))

            if pts:
                gap = Inches(0.14)
                remaining = cards_bottom - top
                total_needed = sum(ch for _, _, ch in pts) + gap * (len(pts) - 1)
                if total_needed > remaining:
                    # 各カードの最小高さは「見出し1行＋desc行数 ÷ 2」程度（descが入り切る最小）
                    # 既存 ch から 0.3インチ引いた値を下限とする（余白を削るが行は削らない）
                    min_chs = [max(Inches(0.62), ch - Inches(0.30)) for _, _, ch in pts]
                    min_total = sum(min_chs) + gap * (len(pts) - 1)
                    if remaining >= min_total:
                        # 入る：必要量を均等に縮める（最小制約付き）
                        # まず最小に縮めて余りを再分配
                        slack = remaining - min_total
                        # 各カードに元のchに比例して余りを配分
                        weights = [ch - mc for (_, _, ch), mc in zip(pts, min_chs)]
                        sum_w = sum(weights)
                        if sum_w > 0:
                            new_chs = [mc + int(slack * w / sum_w) for mc, w in zip(min_chs, weights)]
                        else:
                            new_chs = min_chs
                        pts = [(h, d, nc) for (h, d, _), nc in zip(pts, new_chs)]
                    else:
                        # 全カードを最小にしても入らない：末尾から削減
                        while len(pts) > 1:
                            pts = pts[:-1]
                            min_chs_t = [max(Inches(0.62), ch - Inches(0.30)) for _, _, ch in pts]
                            t_min = sum(min_chs_t) + gap * (len(pts) - 1)
                            if t_min <= remaining:
                                # 縮小して入れる
                                slack = remaining - t_min
                                weights = [ch - mc for (_, _, ch), mc in zip(pts, min_chs_t)]
                                sum_w = sum(weights)
                                if sum_w > 0:
                                    new_chs = [mc + int(slack * w / sum_w) for mc, w in zip(min_chs_t, weights)]
                                else:
                                    new_chs = min_chs_t
                                pts = [(h, d, nc) for (h, d, _), nc in zip(pts, new_chs)]
                                break

                # 描画
                cy = top
                for i, (head, desc, ch) in enumerate(pts):
                    _point_card(s, SAFE, cy, text_w, ch, i + 1, head, desc, accent)
                    cy = cy + ch + gap

        # ─── 右側：画像 or 転スライドの実演プレースホルダ ───
        if img and img_l is not None:
            pic_t = Inches(1.55) + (Inches(0.6) if purpose else Inches(0))
            pic_h = cards_bottom - pic_t
            _card(s, img_l, pic_t, img_w, pic_h, fill=COL_BG_CARD)
            try:
                _add_picture_fit(s, img, img_l + Inches(0.15), pic_t + Inches(0.15),
                                 img_w - Inches(0.3), pic_h - Inches(0.65))
            except Exception:
                pass
            if phase == "転":
                _card(s, img_l, cards_bottom - Inches(0.45), img_w, Inches(0.4), fill=accent)
                _txt(s, img_l, cards_bottom - Inches(0.45), img_w, Inches(0.4),
                     "🎬 実演ポイント（完成後に動画に差し替え可）",
                     size=11, color="FFFFFF", bold=True, font=FONT_HEAD,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            elif visual_role:
                _txt(s, img_l, cards_bottom - Inches(0.35), img_w, Inches(0.3),
                     "💡 " + visual_role,
                     size=10, color=COL_MUTE, font=FONT_BODY, align=PP_ALIGN.CENTER)
        elif phase == "転" and img_l is not None:
            pic_t = Inches(1.55) + (Inches(0.6) if purpose else Inches(0))
            pic_h = cards_bottom - pic_t
            _card(s, img_l, pic_t, img_w, pic_h, fill=tint)
            cy = pic_t + pic_h / 2
            _txt(s, img_l, cy - Inches(0.5), img_w, Inches(0.5),
                 "🎬", size=40, color=accent, bold=True, font=FONT_HEAD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _txt(s, img_l + Inches(0.2), cy, img_w - Inches(0.4), Inches(0.5),
                 "ここに実演写真／動画を入れます",
                 size=13, color=accent, bold=True, font=FONT_HEAD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if visual_role:
                _txt(s, img_l + Inches(0.2), cy + Inches(0.45),
                     img_w - Inches(0.4), Inches(0.5),
                     "想定：" + visual_role,
                     size=10, color=COL_MUTE, font=FONT_BODY,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # =====================================================================
    # 裏表紙
    # =====================================================================
    s = prs.slides.add_slide(blank)
    _set_white_bg(s)
    _txt(s, SAFE, Inches(3.0), SLIDE_W - SAFE * 2, Inches(1.0),
         "ありがとうございました！", size=40, color=COL_INK, bold=True,
         font=FONT_HEAD, align=PP_ALIGN.CENTER)
    _txt(s, SAFE, Inches(4.1), SLIDE_W - SAFE * 2, Inches(0.6),
         content.get("team_line", ""), size=16, color=COL_MUTE,
         font=FONT_BODY, align=PP_ALIGN.CENTER)

    prs.save(output_path)
    return output_path
def _safe_json(text, default=None):
    """LLM出力から最初のJSONオブジェクトを安全に取り出す。```除去にも対応。"""
    if not text:
        return default
    t = text.strip()
    t = re.sub(r"^```(?:json)?", "", t).strip()
    t = re.sub(r"```$", "", t).strip()
    try:
        return json.loads(t)
    except Exception:
        m = re.search(r"\{.*\}", t, re.S)
        if m:
            try:
                return json.loads(m.group(0))
            except Exception:
                pass
    return default


def _first_price(points):
    """価格ポイント群から代表的な金額文字列を1つ取り出す。"""
    for p in points:
        head, desc = _norm_point(p)
        for t in (head, desc):
            if any(k in t for k in ("円", "万", "¥", "価格")):
                return t
    return _norm_point(points[0])[0] if points else ""


def _split(text):
    """説明文を箇条書きへ簡易分割。"""
    if not text:
        return []
    parts = re.split(r"[。\n・]", text)
    return [p.strip() for p in parts if p.strip()][:4]


# ============================================================================
# 3.5) PowerPoint 生成（business調・提案書スタイル）
# ============================================================================
# business方式のカラーパレット（ダークモード／紺背景）
_BIZ_BG          = "0F1B3D"   # 全体背景（深い紺・濃紺）
_BIZ_BG_ELEVATED = "162454"   # 一段明るい紺（タイトル領域の段差）
_BIZ_BG_SURFACE  = "1E3160"   # カード背景（紺系の少し明るい色）
_BIZ_BG_KPI      = "1A2A55"   # KPIパネル背景
_BIZ_NAVY        = "3B82F6"   # アクセント青（明るいブルー、暗背景で映える）
_BIZ_ACCENT      = "F59E0B"   # アクセントオレンジ／ゴールド（暗背景で映える）
_BIZ_TEXT        = "F8FAFC"   # メインテキスト（白に近い）
_BIZ_TEXT_SUB    = "CBD5E1"   # サブテキスト（薄いグレー）
_BIZ_TEXT_MUTE   = "94A3B8"   # 補助テキスト
_BIZ_LINE        = "334155"   # 罫線

# セクション順とラベル（英語/日本語）
_BIZ_SECTION_ORDER = ["problem", "solution", "how", "impact", "why", "next"]


def _biz_section_of(sl):
    s = (sl.get("section") or "").strip().lower()
    return s if s in BIZ_SECTION_LABELS else ""


def build_pptx_biz(content, photos=None, output_path="robot_presentation_biz.pptx",
                   images_by_key=None, cover_image=None):
    """business調（提案書スタイル）のPPTXを生成する。

    セクション順: Problem → Solution → How → Impact → Why → Next（合計7枚＋裏表紙）
    images_by_key: {スライドindex: 画像bytes} （未指定なら photos から impact→how→solution→problem→why→next 順）
    cover_image: 表紙の右側に置く画像（任意）
    """
    slides = content.get("slides", [])

    # 画像の自動配分順
    if images_by_key is None:
        images_by_key = {}
        photos = photos or []
        if photos:
            priority = ["impact", "how", "solution", "problem", "why", "next"]
            ordered = []
            for sec in priority:
                for i, sl in enumerate(slides):
                    if _biz_section_of(sl) == sec and i not in ordered:
                        ordered.append(i)
            for img, i in zip(photos, ordered):
                images_by_key[i] = img
    images_by_key = {k: v for k, v in images_by_key.items() if v}

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # 行高推定ヘルパー
    LINE_H_BODY_IN = 13 * 1.45 / 72
    LINE_H_DESC_IN = 12.5 * 1.25 / 72

    def _chars_per_line(width_emu, font_pt=13):
        in_w = width_emu / 914400.0
        return max(10, int(in_w * 72 / font_pt))

    def _est_lines(text, width_emu, font_pt=13):
        if not text:
            return 0
        cpl = _chars_per_line(width_emu, font_pt)
        lines = 0
        for para in text.split("\n"):
            n = max(1, (len(para) + cpl - 1) // cpl)
            lines += n
        return max(1, lines)

    # =====================================================================
    # 表紙（ダークモード背景）
    # =====================================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s, _BIZ_BG)

    if cover_image:
        hero_l = SLIDE_W - SAFE - Inches(4.6)
        hero_t = Inches(1.0)
        hero_w = Inches(4.6)
        hero_h = Inches(5.5)
        _card(s, hero_l, hero_t, hero_w, hero_h, fill=_BIZ_BG_SURFACE)
        try:
            _add_picture_fit(s, cover_image, hero_l + Inches(0.2), hero_t + Inches(0.2),
                             hero_w - Inches(0.4), hero_h - Inches(0.4))
        except Exception:
            pass
        text_right = hero_l - Inches(0.3)
    else:
        # 表紙左上にアクセント色の細いバンド（端には触れない、安全マージン内）
        band = _card(s, SAFE, Inches(0.7), Inches(2.5), Inches(0.18), fill=_BIZ_ACCENT)
        text_right = SLIDE_W - SAFE

    # 「PROPOSAL」上部ラベル
    _txt(s, SAFE, Inches(1.0), Inches(4.0), Inches(0.4),
         "PROPOSAL  /  事業提案", size=12, color=_BIZ_ACCENT, bold=True,
         font=FONT_HEAD)
    # メインタイトル
    title_w = text_right - SAFE
    _txt(s, SAFE, Inches(1.7), title_w, Inches(2.0),
         content.get("title", ""), size=42, color=_BIZ_TEXT, bold=True, font=FONT_HEAD,
         leading=1.15)
    # サブタイトル
    _txt(s, SAFE, Inches(4.0), title_w, Inches(0.8),
         content.get("subtitle", ""), size=18, color=_BIZ_NAVY, bold=True, font=FONT_HEAD)
    # Executive Summary パネル
    ov = content.get("overview", "")
    if ov:
        _card(s, SAFE, Inches(5.0), title_w, Inches(1.35), fill=_BIZ_BG_SURFACE)
        _txt(s, SAFE + Inches(0.3), Inches(5.08), title_w - Inches(0.6), Inches(0.32),
             "Executive Summary  ｜  本提案の要旨",
             size=11, color=_BIZ_ACCENT, bold=True, font=FONT_HEAD)
        _txt(s, SAFE + Inches(0.3), Inches(5.42), title_w - Inches(0.6), Inches(0.85),
             ov, size=12, color=_BIZ_TEXT, font=FONT_BODY, leading=1.35)
    # チームライン
    _txt(s, SAFE, Inches(6.65), SLIDE_W - SAFE * 2, Inches(0.45),
         content.get("team_line", ""), size=13, color=_BIZ_TEXT_SUB, font=FONT_BODY)

    # =====================================================================
    # 本文（6セクション）
    # =====================================================================
    for idx, sl in enumerate(slides):
        s = prs.slides.add_slide(blank)
        _set_bg(s, _BIZ_BG)
        section = _biz_section_of(sl)
        label_en, label_jp = BIZ_SECTION_LABELS.get(section, ("", ""))

        # ─── ヘッダー：セクション番号＋英語ラベル＋日本語タイトル ───
        # 左に「01」のような大きなセクション番号
        try:
            section_idx = _BIZ_SECTION_ORDER.index(section) + 1
            section_num = f"{section_idx:02d}"
        except ValueError:
            section_num = ""

        if section_num:
            _txt(s, SAFE, Inches(0.45), Inches(0.9), Inches(0.55),
                 section_num, size=32, color=_BIZ_ACCENT, bold=True, font=FONT_HEAD,
                 leading=1.0)
        # 英語小ラベル
        if label_en:
            _txt(s, SAFE + Inches(0.95), Inches(0.5), Inches(5.0), Inches(0.3),
                 label_en.upper(), size=10, color=_BIZ_ACCENT, bold=True, font=FONT_HEAD)
        # タイトル
        title_l = SAFE + Inches(0.95)
        _txt(s, title_l, Inches(0.78), SLIDE_W - title_l - SAFE, Inches(0.55),
             sl.get("title", "") or label_jp, size=24, color=_BIZ_TEXT, bold=True, font=FONT_HEAD)

        # 細い区切り（バーではなく、ヘッダーとボディの境界を示す薄ライン用カード）
        _card(s, SAFE, Inches(1.45), SLIDE_W - SAFE * 2, Pt(1.5), fill=_BIZ_LINE)
        top = Inches(1.65)

        # ─── レイアウト寸法（画像/KPIの有無で分岐） ───
        points = sl.get("points") or []
        kpis = sl.get("kpis") or []
        img = images_by_key.get(idx)
        body_text = sl.get("body", "")
        lead = sl.get("lead", "")
        cta = sl.get("cta", "")
        cards_bottom = Inches(7.15)

        # 画像かKPIブロックを右側に置く（画像優先）
        has_right_panel = bool(img) or bool(kpis)

        if has_right_panel:
            text_w = Inches(5.9)
            right_l = SAFE + text_w + Inches(0.35)
            right_w = SLIDE_W - right_l - SAFE
        else:
            text_w = SLIDE_W - SAFE * 2
            right_l = None
            right_w = None

        # ─── リード文（見出しコピー） ───
        if lead:
            _txt(s, SAFE, top, text_w, Inches(0.45),
                 lead, size=16, color=_BIZ_NAVY, bold=True, font=FONT_HEAD)
            top = top + Inches(0.55)

        # ─── body段落 ───
        if body_text:
            body_lines = _est_lines(body_text, text_w, font_pt=13)
            available_h = cards_bottom - top
            body_h_max = int(available_h * 0.40)
            body_h_needed = Inches(LINE_H_BODY_IN * body_lines + 0.10)
            if body_h_needed > body_h_max:
                body_h_needed = body_h_max
            _txt(s, SAFE, top, text_w, body_h_needed,
                 body_text, size=13, color=_BIZ_TEXT, font=FONT_BODY, leading=1.45)
            top = top + body_h_needed + Inches(0.18)

        # ─── ポイントカード（左テキスト領域内） ───
        if points and top < cards_bottom - Inches(0.4):
            pts = []
            for p in points[:3]:
                if isinstance(p, dict):
                    head = (p.get("head") or "").strip()
                    desc = (p.get("desc") or "").strip()
                else:
                    head, desc = str(p).strip(), ""
                if not head and not desc:
                    continue
                desc_w = text_w - Inches(1.34)
                if desc:
                    desc_lines = _est_lines(desc, desc_w, font_pt=12.5)
                    ch_needed = Inches(0.36 + 0.22 + LINE_H_DESC_IN * desc_lines + 0.16)
                else:
                    ch_needed = Inches(0.65)
                pts.append((head, desc, ch_needed))

            if pts:
                gap = Inches(0.14)
                remaining = cards_bottom - top
                total_needed = sum(ch for _, _, ch in pts) + gap * (len(pts) - 1)
                if total_needed > remaining:
                    min_chs = [max(Inches(0.62), ch - Inches(0.30)) for _, _, ch in pts]
                    min_total = sum(min_chs) + gap * (len(pts) - 1)
                    if remaining >= min_total:
                        slack = remaining - min_total
                        weights = [ch - mc for (_, _, ch), mc in zip(pts, min_chs)]
                        sum_w = sum(weights)
                        if sum_w > 0:
                            new_chs = [mc + int(slack * w / sum_w) for mc, w in zip(min_chs, weights)]
                        else:
                            new_chs = min_chs
                        pts = [(h, d, nc) for (h, d, _), nc in zip(pts, new_chs)]
                    else:
                        while len(pts) > 1:
                            pts = pts[:-1]
                            min_chs_t = [max(Inches(0.62), ch - Inches(0.30)) for _, _, ch in pts]
                            t_min = sum(min_chs_t) + gap * (len(pts) - 1)
                            if t_min <= remaining:
                                slack = remaining - t_min
                                weights = [ch - mc for (_, _, ch), mc in zip(pts, min_chs_t)]
                                sum_w = sum(weights)
                                if sum_w > 0:
                                    new_chs = [mc + int(slack * w / sum_w) for mc, w in zip(min_chs_t, weights)]
                                else:
                                    new_chs = min_chs_t
                                pts = [(h, d, nc) for (h, d, _), nc in zip(pts, new_chs)]
                                break
                # 描画（business調：番号バッジは紺、薄いカード背景）
                cy = top
                for i, (head, desc, ch) in enumerate(pts):
                    _biz_point_card(s, SAFE, cy, text_w, ch, i + 1, head, desc)
                    cy = cy + ch + gap

        # ─── 右パネル：画像 or KPIブロック ───
        if right_l is not None:
            panel_t = Inches(1.65)
            panel_h = cards_bottom - panel_t

            if img:
                # 画像
                _card(s, right_l, panel_t, right_w, panel_h, fill=_BIZ_BG_SURFACE)
                try:
                    _add_picture_fit(s, img, right_l + Inches(0.18), panel_t + Inches(0.18),
                                     right_w - Inches(0.36), panel_h - Inches(0.36))
                except Exception:
                    pass
                # 画像下に小キャプション（visual_role 風だがbizでは控えめ）
                vrole = sl.get("visual_role", "")
                if vrole:
                    _txt(s, right_l, panel_t + panel_h - Inches(0.35),
                         right_w, Inches(0.3),
                         vrole, size=10, color=_BIZ_TEXT_MUTE, font=FONT_BODY,
                         align=PP_ALIGN.CENTER)
            elif kpis:
                # KPIブロック：大きな数字＋ラベル＋出典
                _biz_kpi_panel(s, right_l, panel_t, right_w, panel_h, kpis[:3])

        # ─── 「Next Steps」セクションは下部に CTA バーを置く（端には触れない） ───
        if section == "next" and cta:
            cta_y = cards_bottom + Inches(0.05)
            if cta_y < Inches(7.30):
                _card(s, SAFE, cta_y, SLIDE_W - SAFE * 2, Inches(0.35), fill=_BIZ_ACCENT)
                _txt(s, SAFE, cta_y, SLIDE_W - SAFE * 2, Inches(0.35),
                     "▶  " + cta,
                     size=12, color="FFFFFF", bold=True, font=FONT_HEAD,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # =====================================================================
    # 裏表紙
    # =====================================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s, _BIZ_BG)
    _txt(s, SAFE, Inches(2.8), SLIDE_W - SAFE * 2, Inches(0.5),
         "Thank You", size=44, color=_BIZ_ACCENT, bold=True,
         font=FONT_HEAD, align=PP_ALIGN.CENTER)
    _txt(s, SAFE, Inches(3.7), SLIDE_W - SAFE * 2, Inches(0.5),
         "ご清聴ありがとうございました", size=18, color=_BIZ_TEXT,
         font=FONT_HEAD, align=PP_ALIGN.CENTER)
    _txt(s, SAFE, Inches(4.8), SLIDE_W - SAFE * 2, Inches(0.5),
         content.get("team_line", ""), size=14, color=_BIZ_TEXT_SUB,
         font=FONT_BODY, align=PP_ALIGN.CENTER)

    prs.save(output_path)
    return output_path


def _biz_point_card(slide, l, t, w, h, n, head, desc):
    """business方式の番号付きポイントカード（薄背景、紺の番号バッジ）。"""
    _card(slide, l, t, w, h, fill=_BIZ_BG_SURFACE)
    bd = Inches(0.46)
    bx = l + Inches(0.24)
    if desc:
        by = t + Inches(0.20)
    else:
        by = t + (h - bd) // 2
    _badge(slide, bx, by, bd, str(n), _BIZ_NAVY)
    tx = bx + bd + Inches(0.26)
    tw = l + w - tx - Inches(0.3)
    if desc:
        _txt(slide, tx, t + Inches(0.16), tw, Inches(0.38),
             head, size=15, color=_BIZ_TEXT, bold=True, font=FONT_HEAD, leading=1.1)
        _txt(slide, tx, t + Inches(0.58), tw, h - Inches(0.70),
             desc, size=12, color=_BIZ_TEXT, font=FONT_BODY, leading=1.30)
    else:
        _txt(slide, tx, t, tw, h, head, size=15, color=_BIZ_TEXT, bold=True,
             font=FONT_HEAD, anchor=MSO_ANCHOR.MIDDLE, leading=1.1)


def _biz_kpi_panel(slide, l, t, w, h, kpis):
    """KPI ブロック：大きな数字＋ラベル＋出典を縦に並べる。
    kpis: [{"value": "48歳", "label": "...", "source": "..."}, ...]
    """
    if not kpis:
        return
    # 全体を一段明るい紺のパネルに
    _card(slide, l, t, w, h, fill=_BIZ_BG_KPI)
    # ヘッダー
    _txt(slide, l + Inches(0.25), t + Inches(0.15), w - Inches(0.5), Inches(0.3),
         "Key Metrics  ｜  主要指標",
         size=10, color=_BIZ_ACCENT, bold=True, font=FONT_HEAD)

    # 各KPIカード
    inner_t = t + Inches(0.55)
    inner_h = h - Inches(0.7)
    n = len(kpis)
    gap = Inches(0.12)
    card_h = int((inner_h - gap * (n - 1)) / n) if n > 0 else inner_h
    card_h = max(Inches(1.0), min(card_h, Inches(2.0)))

    cy = inner_t
    for i, k in enumerate(kpis):
        if cy + card_h > t + h - Inches(0.1):
            break
        # 個別カード（さらに一段明るい紺：surface色でコントラストを出す）
        _card(slide, l + Inches(0.2), cy, w - Inches(0.4), card_h, fill=_BIZ_BG_SURFACE)
        # 大きな数字（アクセント色）
        val = str(k.get("value", "")).strip()
        lbl = str(k.get("label", "")).strip()
        src = str(k.get("source", "")).strip()
        if val:
            _txt(slide, l + Inches(0.35), cy + Inches(0.08),
                 w - Inches(0.7), Inches(0.55),
                 val, size=26, color=_BIZ_ACCENT, bold=True, font=FONT_HEAD,
                 leading=1.0)
        if lbl:
            _txt(slide, l + Inches(0.35), cy + Inches(0.65),
                 w - Inches(0.7), Inches(0.3),
                 lbl, size=11, color=_BIZ_TEXT, font=FONT_BODY, leading=1.2)
        if src:
            _txt(slide, l + Inches(0.35), cy + card_h - Inches(0.28),
                 w - Inches(0.7), Inches(0.22),
                 "出典: " + src, size=9, color=_BIZ_TEXT_MUTE, font=FONT_BODY)
        cy = cy + card_h + gap


def pdf_first_page_png(pdf_bytes, dpi=200):
    """PDFの1ページ目をPNG bytesにして返す（OCR入力用）。pymupdf必須。"""
    import fitz
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    page = doc[0]
    pix = page.get_pixmap(dpi=dpi)
    return pix.tobytes("png")


# ============================================================================
# ベータ版PowerPointの取り込み
# ============================================================================
def pptx_extract_media_images(pptx_bytes, min_kb=12):
    """ベータ版PPTに埋め込まれた写真（jpg/png等）だけを取り出して bytes のリストで返す。
    スライド全体の画像化ではなく“中の写真”を拾うので、実機写真などの素材抽出に向く。
    min_kb 未満（アイコン等）は除外。"""
    import zipfile
    out = []
    seen = set()
    with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as z:
        names = [n for n in z.namelist() if n.startswith("ppt/media/")]
        names.sort(key=_natural_key)
        for n in names:
            ext = os.path.splitext(n)[1].lower()
            if ext not in (".jpg", ".jpeg", ".png", ".bmp", ".gif", ".tiff"):
                continue
            data = z.read(n)
            if len(data) < min_kb * 1024:        # 小さすぎる画像（装飾・アイコン）は除外
                continue
            key = hashlib.md5(data).hexdigest()
            if key in seen:
                continue
            seen.add(key)
            out.append(data)
    return out


def pptx_extract_text(pptx_bytes):
    """各スライドのテキストを取り出して [{'slide':n, 'text':...}] を返す（python-pptxのみ・常に動作）。"""
    from pptx import Presentation as _P
    prs = _P(io.BytesIO(pptx_bytes))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        chunks = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                t = "\n".join(p.text for p in sh.text_frame.paragraphs if p.text.strip())
                if t.strip():
                    chunks.append(t.strip())
        out.append({"slide": i, "text": "\n".join(chunks)})
    return out


def _natural_key(name):
    m = re.search(r"(\d+)", os.path.basename(name))
    return int(m.group(1)) if m else 0


def pptx_to_slide_images(pptx_bytes):
    """ベータ版PPTの各スライドをPNG bytesのリストにして返す。
    Windowsの PowerPoint(COM) を優先し、無ければ LibreOffice(soffice) を使用。
    どちらも無ければ RuntimeError。"""
    import tempfile
    tmpdir = tempfile.mkdtemp(prefix="draftppt_")
    src = os.path.join(tmpdir, "draft.pptx")
    with open(src, "wb") as f:
        f.write(pptx_bytes)

    # 1) PowerPoint COM（Windows・PowerPointインストール済みなら最も確実）
    try:
        return _pptx_images_powerpoint(src, tmpdir)
    except Exception:
        pass
    # 2) LibreOffice 経由（soffice + pdftoppm が必要）
    try:
        return _pptx_images_soffice(src, tmpdir)
    except Exception as e:
        raise RuntimeError(
            "スライド画像化には PowerPoint(COM) または LibreOffice が必要です。"
            f"（詳細: {e}）画面キャプチャの貼り付けもご利用ください。")


def _pptx_images_powerpoint(src, tmpdir):
    import glob
    import comtypes.client
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    try:
        powerpoint.Visible = 1
    except Exception:
        pass
    pres = powerpoint.Presentations.Open(src, WithWindow=False)
    out_dir = os.path.join(tmpdir, "png")
    os.makedirs(out_dir, exist_ok=True)
    try:
        pres.Export(out_dir, "PNG")  # 全スライドをPNG出力（言語によりファイル名が異なる）
    finally:
        pres.Close()
        powerpoint.Quit()
    files = sorted(glob.glob(os.path.join(out_dir, "*.PNG")) +
                   glob.glob(os.path.join(out_dir, "*.png")), key=_natural_key)
    return [open(p, "rb").read() for p in files]


def _pptx_images_soffice(src, tmpdir):
    import glob
    import subprocess
    import shutil
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise RuntimeError("soffice not found")
    subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", tmpdir, src],
                   check=True, timeout=120,
                   stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    pdf = os.path.join(tmpdir, "draft.pdf")
    pdftoppm = shutil.which("pdftoppm")
    if pdftoppm:
        subprocess.run([pdftoppm, "-png", "-r", "150", pdf, os.path.join(tmpdir, "slide")],
                       check=True, timeout=120)
        files = sorted(glob.glob(os.path.join(tmpdir, "slide-*.png")), key=_natural_key)
        return [open(p, "rb").read() for p in files]
    # pdftoppm が無ければ pymupdf でPDF→PNG
    import fitz
    doc = fitz.open(pdf)
    return [doc[i].get_pixmap(dpi=150).tobytes("png") for i in range(len(doc))]


# ============================================================================
# 8) プロジェクトIO（編集中状態＋生成済みPPTを単一の.zipで保存／復元／逆輸入）
# ============================================================================
PROJECT_VERSION = 1
PROJECT_MANIFEST = "project.json"


def _b64_encode(b):
    """画像bytesをbase64文字列へ。jsonに埋め込まずファイル分離するので将来用に残す。"""
    import base64
    return base64.b64encode(b).decode("ascii") if b else ""


def _b64_decode(s):
    import base64
    return base64.b64decode(s.encode("ascii")) if s else b""


def save_project_zip(state, output_path=None):
    """編集中状態と生成物を1つの.zipにまとめて返す。

    引数:
      state: dict（Streamlit側の session_state から必要分を切り出したもの）
        必須/任意キー:
          - "idea": dict           OCR後の編集中アイデアシート
          - "sheet_img": bytes     アイデアシート画像（PNG）  ※任意
          - "sheet_mime": str      "image/png" など
          - "pasted_photos": [bytes, ...]   貼り付けで追加した写真
          - "draft_slides": [bytes, ...]    ベータ版PPTから取り込んだスライド画像
          - "draft_text":   [str, ...]      同テキスト
          - "uploaded_files_data": [bytes, ...] ファイルアップロードした写真
          - "cover_image": bytes   表紙アイコン  ※任意
          - "images_by_key":      {str: bytes}    従来方式の画像割り当て
          - "sato_images_by_key": {int: bytes}    佐藤方式の画像割り当て
          - "normal_pptx_bytes": bytes  生成済み従来方式PPT  ※任意
          - "sato_pptx_bytes":   bytes  生成済み佐藤方式PPT  ※任意
          - "normal_content": dict  生成原稿  ※任意
          - "sato_content":   dict  生成原稿  ※任意
          - "sato_research": str    Gemini調査結果  ※任意
          - "robot_name": str       ファイル名生成用

      output_path: 出力パス。None なら BytesIO を返す。

    返り値:
      output_path 指定時はそのパス、None時は zip の bytes。
    """
    import zipfile

    state = state or {}
    robot_name = (state.get("robot_name") or "robot").strip() or "robot"
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

    # ---- バイナリファイルをzip内パスへ並べる ----
    file_map = {}   # zip内パス → bytes

    def _put_bytes(zip_path, data):
        if data:
            file_map[zip_path] = data

    def _put_list(prefix, items):
        # items は bytes リスト、何番目かをファイル名に
        paths = []
        for i, b in enumerate(items or []):
            zp = f"{prefix}/{i:03d}.png"
            file_map[zp] = b
            paths.append(zp)
        return paths

    # 各種ファイル
    sheet_path = None
    if state.get("sheet_img"):
        ext = ".pdf" if state.get("sheet_mime") == "application/pdf" else ".png"
        sheet_path = f"sheet/sheet{ext}"
        file_map[sheet_path] = state["sheet_img"]

    pasted_paths = _put_list("pasted", state.get("pasted_photos", []))
    draft_slide_paths = _put_list("draft_slides", state.get("draft_slides", []))
    uploaded_paths = _put_list("uploaded", state.get("uploaded_files_data", []))

    cover_path = None
    if state.get("cover_image"):
        cover_path = "cover/cover.png"
        file_map[cover_path] = state["cover_image"]

    # 画像割り当て（key → 画像bytes）→ ファイル分離
    images_by_key_paths = {}
    for k, b in (state.get("images_by_key") or {}).items():
        if not b:
            continue
        # キーをファイル名に使えるよう正規化
        safe_k = re.sub(r"[^A-Za-z0-9_-]", "_", str(k))
        zp = f"assign_normal/{safe_k}.png"
        file_map[zp] = b
        images_by_key_paths[str(k)] = zp

    sato_images_by_key_paths = {}
    for k, b in (state.get("sato_images_by_key") or {}).items():
        if not b:
            continue
        zp = f"assign_sato/{int(k):02d}.png"
        file_map[zp] = b
        sato_images_by_key_paths[str(int(k))] = zp

    biz_images_by_key_paths = {}
    for k, b in (state.get("biz_images_by_key") or {}).items():
        if not b:
            continue
        zp = f"assign_biz/{int(k):02d}.png"
        file_map[zp] = b
        biz_images_by_key_paths[str(int(k))] = zp

    # 生成済みPPT
    normal_pptx_path = None
    if state.get("normal_pptx_bytes"):
        normal_pptx_path = "output/normal.pptx"
        file_map[normal_pptx_path] = state["normal_pptx_bytes"]

    sato_pptx_path = None
    if state.get("sato_pptx_bytes"):
        sato_pptx_path = "output/sato.pptx"
        file_map[sato_pptx_path] = state["sato_pptx_bytes"]

    biz_pptx_path = None
    if state.get("biz_pptx_bytes"):
        biz_pptx_path = "output/biz.pptx"
        file_map[biz_pptx_path] = state["biz_pptx_bytes"]

    # ---- マニフェスト ----
    manifest = {
        "version": PROJECT_VERSION,
        "saved_at": timestamp,
        "robot_name": robot_name,
        "idea": state.get("idea") or {},
        "sheet_mime": state.get("sheet_mime", "image/png"),
        "draft_text": state.get("draft_text", []),
        # パス参照
        "paths": {
            "sheet": sheet_path,
            "cover": cover_path,
            "pasted_photos": pasted_paths,
            "draft_slides": draft_slide_paths,
            "uploaded_files": uploaded_paths,
            "images_by_key": images_by_key_paths,
            "sato_images_by_key": sato_images_by_key_paths,
            "biz_images_by_key": biz_images_by_key_paths,
            "normal_pptx": normal_pptx_path,
            "sato_pptx": sato_pptx_path,
            "biz_pptx": biz_pptx_path,
        },
        "normal_content": state.get("normal_content"),
        "sato_content": state.get("sato_content"),
        "biz_content": state.get("biz_content"),
        "sato_research": state.get("sato_research", ""),
        "biz_research": state.get("biz_research", ""),
    }

    # ---- zip書き出し ----
    if output_path is None:
        buf = io.BytesIO()
        zip_target = buf
    else:
        zip_target = output_path

    with zipfile.ZipFile(zip_target, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(PROJECT_MANIFEST,
                    json.dumps(manifest, ensure_ascii=False, indent=2))
        for zp, data in file_map.items():
            zf.writestr(zp, data)

    if output_path is None:
        return buf.getvalue()
    return output_path


def load_project_zip(zip_bytes_or_path):
    """save_project_zip で書き出した.zipを読み込み、復元用 state dict を返す。

    返す dict は save_project_zip の引数とほぼ同じ構造（bytesが復元されている）。
    """
    import zipfile

    if isinstance(zip_bytes_or_path, (bytes, bytearray)):
        z_in = io.BytesIO(zip_bytes_or_path)
    else:
        z_in = zip_bytes_or_path

    with zipfile.ZipFile(z_in, "r") as zf:
        manifest = json.loads(zf.read(PROJECT_MANIFEST).decode("utf-8"))

        def _read(zp):
            if not zp:
                return None
            try:
                return zf.read(zp)
            except KeyError:
                return None

        def _read_list(zp_list):
            return [zf.read(zp) for zp in (zp_list or []) if zp in zf.namelist()]

        paths = manifest.get("paths", {})

        # 画像割り当て：従来方式
        images_by_key = {}
        for k, zp in (paths.get("images_by_key") or {}).items():
            b = _read(zp)
            if b:
                images_by_key[k] = b

        # 画像割り当て：佐藤方式（キーはint）
        sato_images_by_key = {}
        for k, zp in (paths.get("sato_images_by_key") or {}).items():
            b = _read(zp)
            if b:
                try:
                    sato_images_by_key[int(k)] = b
                except ValueError:
                    sato_images_by_key[k] = b

        # 画像割り当て：business方式（キーはint）
        biz_images_by_key = {}
        for k, zp in (paths.get("biz_images_by_key") or {}).items():
            b = _read(zp)
            if b:
                try:
                    biz_images_by_key[int(k)] = b
                except ValueError:
                    biz_images_by_key[k] = b

        state = {
            "version": manifest.get("version", 1),
            "saved_at": manifest.get("saved_at"),
            "robot_name": manifest.get("robot_name", ""),
            "idea": manifest.get("idea", {}),
            "sheet_mime": manifest.get("sheet_mime", "image/png"),
            "sheet_img": _read(paths.get("sheet")),
            "cover_image": _read(paths.get("cover")),
            "pasted_photos": _read_list(paths.get("pasted_photos")),
            "draft_slides": _read_list(paths.get("draft_slides")),
            "draft_text": manifest.get("draft_text", []),
            "uploaded_files_data": _read_list(paths.get("uploaded_files")),
            "images_by_key": images_by_key,
            "sato_images_by_key": sato_images_by_key,
            "biz_images_by_key": biz_images_by_key,
            "normal_pptx_bytes": _read(paths.get("normal_pptx")),
            "sato_pptx_bytes": _read(paths.get("sato_pptx")),
            "biz_pptx_bytes": _read(paths.get("biz_pptx")),
            "normal_content": manifest.get("normal_content"),
            "sato_content": manifest.get("sato_content"),
            "biz_content": manifest.get("biz_content"),
            "sato_research": manifest.get("sato_research", ""),
            "biz_research": manifest.get("biz_research", ""),
        }
    return state


def reimport_external_pptx(state, pptx_bytes, kind="sato"):
    """ユーザーがPowerPoint側で手編集した.pptxをプロジェクトに取り込む（逆輸入）。

    kind: "sato" / "normal" / "biz" のいずれか。指定方式の生成済みPPTを上書きする。
    state: load_project_zip 互換の dict（in-placeで更新）または None。
    返り値: 更新後の state dict
    """
    state = dict(state or {})
    if kind == "sato":
        state["sato_pptx_bytes"] = pptx_bytes
    elif kind == "normal":
        state["normal_pptx_bytes"] = pptx_bytes
    elif kind == "biz":
        state["biz_pptx_bytes"] = pptx_bytes
    else:
        raise ValueError(f"unknown kind: {kind}")
    return state

